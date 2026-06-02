#!/usr/bin/env python3
"""Generate KmpSDK stakeholder PowerPoint. Requires: pip install python-pptx"""

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    raise SystemExit("Install dependency: pip install python-pptx")

OUTPUT = Path(__file__).resolve().parents[1] / "docs" / "presentations" / "KmpSDK-Stakeholder-Overview.pptx"

BLUE = RGBColor(0x1A, 0x56, 0xDB)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_title(slide, text, subtitle=None):
    if slide.shapes.title:
        slide.shapes.title.text = text
        for p in slide.shapes.title.text_frame.paragraphs:
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = DARK
    if subtitle and len(slide.placeholders) > 1:
        ph = slide.placeholders[1]
        ph.text = subtitle
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18)
            p.font.color.rgb = GRAY


def add_bullets(slide, items, left=0.8, top=1.6, width=8.5, height=5.0, size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)


def add_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = [
        (
            "title",
            "KmpSDK",
            "Headless Kotlin Multiplatform SDK for Android & iOS\nPersonal open-source project · Maven Central · v1.0.0",
            "Open with: who you are, that this is a personal initiative (not a company product), and today's goal — awareness + optional adoption.",
        ),
        (
            "content",
            "Agenda",
            None,
            [
                "Problem we solve on mobile projects",
                "What KmpSDK is (and is not)",
                "Architecture at a glance",
                "Three integration paths — pick per feature",
                "Business value for teams",
                "Availability, governance & next steps",
                "Q&A",
            ],
            "Keep to 20–25 minutes total; leave 5 min for questions.",
        ),
        (
            "content",
            "The challenge on mobile projects",
            None,
            [
                "Android and iOS often duplicate networking, auth, caching, and offline logic",
                "Each new feature repeats boilerplate: API client, error handling, sync, state management",
                "UI frameworks differ — but infrastructure problems are the same",
                "Result: slower delivery, inconsistent patterns, harder maintenance",
            ],
            "Stakeholders care about time-to-market and risk — frame pain in business terms.",
        ),
        (
            "content",
            "What is KmpSDK?",
            None,
            [
                "A shared Kotlin Multiplatform library — one codebase for Android & iOS infrastructure",
                "Headless: provides contracts & services, not UI widgets",
                "Host apps keep full control of design (Compose, SwiftUI, XML, etc.)",
                "Published on Maven Central — any KMP team can add one dependency line",
                "Personal project by Nitesh Kukreja — open source on GitHub",
            ],
            "Emphasize: SDK ≠ finished app. It accelerates app teams.",
        ),
        (
            "content",
            "What KmpSDK is NOT",
            None,
            [
                "Not a replacement for product or UX teams",
                "Not a no-code platform — developers still build screens and business rules",
                "Not a company-mandated product (unless leadership chooses to adopt it)",
                "Not limited to one app — reusable across Kotlin Multiplatform projects",
            ],
            "Sets expectations for non-engineering stakeholders.",
        ),
        (
            "content",
            "Architecture — high level",
            None,
            [
                "Host app layer: UI, ViewModels, navigation (your brand & flows)",
                "KmpSDK shared layer: networking, auth, logging, optional offline sync, MVI helpers",
                "Platform layer: Android (OkHttp) and iOS (Darwin) drivers",
                "Flow: Maven dependency → KmpSdk.init() → feature modules → your UI",
            ],
            "Point to diagram image on next slide if you embedded kmpsdk-architecture-flow.png.",
        ),
        (
            "content",
            "What the SDK provides vs what teams build",
            None,
            [
                "SDK provides: HTTP client, auth plugin, cache/queue, sync coordinator, MVI base types",
                "Teams build: DTOs, use cases, ViewModels, all visual design",
                "Optional depth: full offline-first only when the product needs it",
                "Flexible: simple screens need less code than offline-heavy catalogs",
            ],
            "Use the README comparison table if asked for detail.",
        ),
        (
            "content",
            "Three integration paths (per feature)",
            None,
            [
                "Path A — Online only: login, forms, settings (minimal code, no app database)",
                "Path B — Network-first + SDK cache: show last API response when offline",
                "Path C — Full offline-first: catalogs, field apps (SQL + sync in your app)",
                "Teams choose per screen — not one-size-fits-all",
            ],
            "Decision: Need data in YOUR database offline? Yes → C. No → A or B.",
        ),
        (
            "content",
            "Path A — Online only (example)",
            None,
            [
                "Use case calls KmpSdk.networkClient.get/post",
                "No local database tables for that feature",
                "Fastest path for simple APIs",
                "Good fit: auth, profile update, one-shot workflows",
            ],
            "Technical audience: mention NETWORK_FIRST init flags.",
        ),
        (
            "content",
            "Path C — Offline-first (example)",
            None,
            [
                "Your SQLDelight tables + local + remote + repository pattern",
                "SDK supplies BaseSyncRepository, bindSyncList, dirty sync, offline queue",
                "Good fit: product lists, orders, inspections with poor connectivity",
                "Feature generator CLI scaffolds boilerplate for Path C",
            ],
            "Only teams with real offline requirements invest in Path C.",
        ),
        (
            "content",
            "Key capabilities (v1.0 / v1.4)",
            None,
            [
                "Authentication & token refresh",
                "Structured errors & logging (redacted headers)",
                "HTTP cache & offline mutation queue (configurable)",
                "Multi-environment config (dev / staging / prod)",
                "Telemetry hooks, remote config, certificate pinning",
                "REST helpers & feature generator for faster delivery",
            ],
            "Pick 2–3 bullets most relevant to your org (e.g. offline, security, multi-env).",
        ),
        (
            "content",
            "Business value",
            None,
            [
                "Faster feature delivery — shared patterns across Android & iOS",
                "Lower duplication cost — one investment in infrastructure",
                "Consistent quality — auth, errors, sync behave the same everywhere",
                "Easier onboarding — documented paths A/B/C and step-by-step README",
                "Open source — no vendor lock-in; inspect and contribute on GitHub",
            ],
            "Translate to KPIs your org uses: cycle time, defect rate, platform parity.",
        ),
        (
            "content",
            "Who benefits?",
            None,
            [
                "Mobile engineers — less boilerplate, clear architecture",
                "Engineering managers — predictable structure across squads",
                "QA — consistent error/offline behaviour to test against",
                "Architecture / platform teams — reference implementation for KMP",
                "Other departments — faster mobile delivery for digital initiatives",
            ],
            "Invite stakeholders to name a pilot app or squad.",
        ),
        (
            "content",
            "Availability & trust",
            None,
            [
                "Maven Central: in.co.niteshkukreja:kmp-sdk:1.0.0",
                "GitHub: github.com/NiteshKuk/Kotlin-Multiplatform-SDK",
                "Apache 2.0 license",
                "Documented integration guide + contributing guidelines",
                "Branch protection & controlled releases (quality gate)",
            ],
            "Non-engineers: Maven Central = standard public registry for dependencies.",
        ),
        (
            "content",
            "Governance (personal project)",
            None,
            [
                "Personal initiative — not automatically an enterprise standard",
                "Collaborators contribute via pull request; owner reviews merges",
                "Only owner publishes new versions to Maven Central",
                "Organizations can fork, audit, or adopt with their own approval process",
            ],
            "Important for legal/compliance: clarify adoption path inside your company.",
        ),
        (
            "content",
            "Suggested next steps",
            None,
            [
                "Share GitHub README with mobile chapter leads",
                "Identify one pilot feature (Path A or C) on an existing KMP app",
                "30-minute technical demo for interested developers",
                "Gather feedback — roadmap driven by real team needs",
                "Optional: internal guild session on Kotlin Multiplatform",
            ],
            "End with a clear ask: pilot squad, demo date, or sponsorship.",
        ),
        (
            "content",
            "Thank you · Q&A",
            None,
            [
                "GitHub: github.com/NiteshKuk/Kotlin-Multiplatform-SDK",
                "Dependency: implementation(\"in.co.niteshkukreja:kmp-sdk:1.0.0\")",
                "Contact: [add your email / Slack / LinkedIn]",
                "Questions?",
            ],
            "Have README and architecture diagram ready on second screen.",
        ),
    ]

    for item in slides_data:
        kind = item[0]
        title = item[1]
        subtitle = item[2]
        bullets = item[3] if len(item) > 3 else None
        notes = item[4] if len(item) > 4 else ""

        if kind == "title":
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            set_title(slide, title, subtitle)
        else:
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            set_title(slide, title)
            if bullets:
                add_bullets(slide, bullets)
        add_notes(slide, notes)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Created: {OUTPUT}")


if __name__ == "__main__":
    build()
